from pptx import Presentation
from pptx.util import Inches
import io
from services.shared.chunker import chunk_by_sections
from typing import List
import logging

logger = logging.getLogger(__name__)


def parse_pptx(file_bytes: bytes, filename: str, source_id: str) -> List[dict]:
    """
    Extract text from each slide of a PPTX file.
    Preserves slide number and slide title for citation.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    sections = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        slide_title = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Identify title shape
            if shape.shape_type == 13 or (hasattr(shape, "name") and "title" in shape.name.lower()):
                slide_title = shape.text_frame.text.strip()

            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    slide_texts.append(line)

        text = "\n".join(slide_texts).strip()
        if not text:
            continue

        sections.append({
            "text": text,
            "source_type": "pptx",
            "source_id": source_id,
            "source_label": filename,
            "slide": slide_num,
            "slide_title": slide_title or f"Slide {slide_num}",
            "total_slides": len(prs.slides),
        })

    if not sections:
        logger.warning(f"No text extracted from PPTX: {filename}")
        return []

    chunks = chunk_by_sections(sections)
    logger.info(f"PPTX '{filename}': {len(chunks)} chunks from {len(sections)} slides")
    return chunks


def summarize_pptx_metadata(file_bytes: bytes, filename: str) -> dict:
    prs = Presentation(io.BytesIO(file_bytes))
    return {
        "filename": filename,
        "slides": len(prs.slides),
        "title": filename,
    }
